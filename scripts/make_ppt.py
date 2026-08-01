"""Generate the basic (text-only) NutriBasket final-year project presentation.

Usage: python scripts/make_ppt.py
Output: NUTRIBASKET_PROJECT_PRESENTATION.pptx in the repo root.
"""

from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT = "NUTRIBASKET_PROJECT_PRESENTATION.pptx"

SLIDES = [
    ("NutriBasket", [
        "A Barcode-Scanning Nutrition Assistant for Smart Grocery Shopping",
        "",
        "Final Year Project (B.Tech)",
        "Submitted by: [Student 1 Name], [Student 2 Name], [Student 3 Name]",
        "Under the guidance of: [Guide Name, Designation]",
        "[Department Name], [College / Institute Name]",
        "Academic Year 20XX - 20XX",
    ]),
    ("Introduction & Problem Statement", [
        "Scan a product barcode with your phone camera and instantly see its nutrition info: calories, protein, carbs, fat.",
        "Data comes from Open Food Facts - a free, open, crowd-sourced database of packaged foods.",
        "Problem: nutrition labels are dense and hard to compare; barcode numbers carry no human-readable meaning.",
        "Existing grocery apps are ordering-focused; nutrition-tracking apps require tedious manual entry.",
        "NutriBasket combines camera scanning, nutrition lookup, and a priced basket in one flow.",
    ]),
    ("Objectives", [
        "1. Camera-based barcode and QR scanner that works in mobile browsers.",
        "2. Manual barcode entry fallback with GTIN check-digit validation.",
        "3. Backend service fetching and normalizing data from the Open Food Facts API.",
        "4. Per-product nutrition display with user-entered prices (INR).",
        "5. Nutrition-aware basket with running totals through a simulated checkout flow.",
        "6. End-to-end deployment with automated pipelines, verified against live data.",
    ]),
    ("System Architecture", [
        "Three-tier web architecture:",
        "- Frontend (React 19 + Vite 8) deployed on Vercel - scanner, product, basket, checkout screens.",
        "- Backend (Node.js + Express) deployed on Render - validation, caching, Open Food Facts lookup.",
        "- Open Food Facts API - the only external data source.",
        "Backend resilience: 10 s timeout, retry with backoff, 24-hour in-memory cache.",
        "Deployment: GitHub Actions workflow triggers Render deploys; Vercel rebuilds on push.",
    ]),
    ("Technology Stack", [
        "Frontend: React 19, Vite 8, @zxing/browser (barcode decoding), plain CSS (emerald theme).",
        "Backend: Node.js, Express 4, cors, helmet, dotenv.",
        "Data: Open Food Facts API v2 (no API key required), in-memory cache (no database).",
        "Hosting: Vercel (frontend), Render free tier (backend), GitHub Actions (auto-deploy).",
        "Version control: Git, fork-based workflow with pull requests to upstream.",
    ]),
    ("Implementation Highlights", [
        "Real camera scanning with ZXing: permission denied / no camera / busy camera states with retry.",
        "Manual entry fallback validated by GTIN check-digit algorithm on both client and server.",
        "Normalized API responses: barcode, name, brand, image, pack quantity, nutrition per 100 g / serving.",
        "Price input guarded with /^\\d+(\\.\\d{1,2})?$/ so invalid values can never reach basket totals (NaN bug fixed).",
        "Error boundary prevents white-screen crashes; cold-start hint for the free-tier backend.",
        "Seven-screen flow: welcome, home, scan, product, basket, checkout, success.",
    ]),
    ("Testing & Results", [
        "Manual, verification-driven testing: linting (0 errors), production builds, live API checks, browser-level testing.",
        "T1 valid barcode (Nutella 3017624010701) -> 200 normalized product - PASS.",
        "T2 unknown barcode -> 404 not-found state - PASS.",
        "T3 invalid barcode -> 400 invalid_barcode, blocked client-side - PASS.",
        "T4 price validation rejects e, 1e5, 1.2.3, -5 - PASS.",
        "T5-T7 camera permission denied / no camera / cold start handled with hints - PASS.",
        "T8-T10 health, service root, basket totals - PASS.",
    ]),
    ("Verification against Live Deployment", [
        "Frontend: https://nutri-basket-six.vercel.app - 200.",
        "Backend: https://nutribasket-api.onrender.com - service info at /, health OK at /health.",
        "Product lookup 3017624010701 -> 200; unknown barcode -> 404.",
        "CORS restricted to the frontend origin via CLIENT_ORIGIN.",
        "Every push auto-deploys via GitHub Actions (verified live).",
    ]),
    ("Challenges & Solutions", [
        "NaN poisoning of basket totals -> strict price-input validation with a regex guard.",
        "Free-tier cold start (~40 s after 15 min idle) -> user-visible hint after 5 s.",
        "Camera quirks (iOS autoplay, StrictMode double-mount) -> start on user tap, safe stream cleanup.",
        "Unknown OFF codes returning status 1 -> server guards on product_name for real 404s.",
        "Secrets hygiene -> API keys kept out of the repo; .gitignore blocks secret file patterns.",
    ]),
    ("Future Scope", [
        "Persistence and accounts: PostgreSQL/Supabase for baskets, prices, scan history.",
        "Automated tests: unit tests for GTIN validation and API client; end-to-end flow tests.",
        "State management migration to Zustand as screens multiply.",
        "Real UPI payment gateway instead of the simulated QR screen.",
        "Performance: lazy-load the scanner to cut the ~687 KB initial bundle.",
        "AI-based nutrition scoring and recommendations; torch / camera-toggle controls.",
    ]),
    ("Conclusion", [
        "All six objectives were achieved and verified end-to-end against live data.",
        "Thin frontend + resilient caching backend keeps the system simple, free to run, and easy to reproduce.",
        "Open data, browser camera APIs, and free hosting tiers combine into a practical nutrition-shopping tool.",
        "The architecture and code are open-source and self-deployable from the repository.",
    ]),
    ("Thank You", [
        "Questions welcome.",
        "",
        "Live demo: https://nutri-basket-six.vercel.app",
        "Repository: https://github.com/Dutta-aman/NutriBasket",
    ]),
]


def add_title_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = ""
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        if i == 1:
            p.font.size = Pt(20)
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.3), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for idx, (title, bullets) in enumerate(SLIDES):
        if idx == 0:
            add_title_slide(prs, title, bullets)
        else:
            add_content_slide(prs, title, bullets)
    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
